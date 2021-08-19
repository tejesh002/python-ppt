from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import pandas as pd

def python_ppt(excel_file,output_file):
    df = pd.read_excel(excel_file)
    root = Presentation()
    for index, row in df.iterrows():
        first_slide_layout = root.slide_layouts[5]
        slide = root.slides.add_slide(first_slide_layout)
        slide.shapes.title.text = row["Name"]

        txBox= slide.shapes.add_textbox(Inches(3.5),Inches(2),Inches(1),Inches(1))
        tf = txBox.text_frame
        tf.text = "Description 1"
        p = tf.add_paragraph() 
        p.text = row['Description 1']

        txBox1= slide.shapes.add_textbox(Inches(6.5),Inches(2),Inches(1),Inches(1))
        tf1 = txBox1.text_frame
        tf1.text = "Description 2"
        p1 = tf1.add_paragraph()
        p1.text = row['Description 2']

        chart_data = ChartData()
        chart_data1 = ChartData()

        chart_data.categories = ["Val1", "Val2"]  

        chart_data.add_series("",
                            (int(row["Val 1 "]), 
                                int(row["Val 2"]))) 

        chart_data1.categories = ['Val3', 'Val4']  

        chart_data1.add_series("",
                            (int(row["Val 3"]), 
                                int(row["Val 4"]))) 
        
        x, y, cx, cy = Inches(3), Inches(3), Inches(2), Inches(2) 

        x1,y1,cx1,cy1 = Inches(6),Inches(3), Inches(2), Inches(2)
        
        slide.shapes.add_chart( XL_CHART_TYPE.PIE, x,
                            y, cx, cy, chart_data )
        slide.shapes.add_chart( XL_CHART_TYPE.PIE, x1,
                            y1, cx1, cy1, chart_data1 )
        

    root.save(output_file)


python_ppt("Problem Statement.xlsx","output.pptx")